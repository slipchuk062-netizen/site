from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
COLOR_BG = RGBColor(26, 26, 26)  # Dark background
COLOR_GREEN = RGBColor(16, 185, 129)  # Main green
COLOR_DARK_GREY = RGBColor(60, 60, 60)  # Dark grey for boxes
COLOR_LIGHT_GREY = RGBColor(200, 200, 200)  # Light grey text
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LIGHT_GREEN_BG = RGBColor(220, 252, 231)  # Light green for solution boxes

def add_bg_rectangle(slide, color=COLOR_BG):
    """Add dark background to slide"""
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)

def add_rounded_box(slide, left, top, width, height, text, bg_color, text_color=COLOR_WHITE, font_size=14):
    """Add a rounded rectangle box with text"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = COLOR_GREEN
    shape.line.width = Pt(2)
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = 'Arial'
    
    return shape

def add_header(slide, text):
    """Add header text box"""
    shape = add_rounded_box(
        slide,
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.6),
        text, COLOR_GREEN, COLOR_WHITE, 28
    )
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_slide_number(slide, num, total):
    """Add slide number at bottom"""
    textbox = slide.shapes.add_textbox(
        Inches(8.5), Inches(6.9), Inches(1), Inches(0.3)
    )
    p = textbox.text_frame.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.RIGHT

# SLIDE 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rectangle(slide)

# University info
textbox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "Міністерство освіти і науки України\n[Ваш Навчальний Заклад]"
p.font.size = Pt(16)
p.font.color.rgb = COLOR_LIGHT_GREY
p.alignment = PP_ALIGN.CENTER
p.line_spacing = 1.5

# Main title
title_box = add_rounded_box(
    slide,
    Inches(1), Inches(2),
    Inches(8), Inches(1.5),
    "Інтелектуальна платформа туристичної аналітики\nЖитомирської області",
    COLOR_DARK_GREY, COLOR_GREEN, 32
)
title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_box.text_frame.paragraphs[0].font.bold = True

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.7), Inches(7), Inches(0.5))
p = subtitle_box.text_frame.paragraphs[0]
p.text = "з використанням кластеризації та геопросторового аналізу"
p.font.size = Pt(18)
p.font.color.rgb = COLOR_LIGHT_GREY
p.alignment = PP_ALIGN.CENTER

# Author info
author_box = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(6), Inches(1.5))
tf = author_box.text_frame
p = tf.paragraphs[0]
p.text = "[Ваше Прізвище, Ім'я, По батькові]\n122 «Комп'ютерні науки»\n\nКерівник роботи:\n[ПІБ Керівника, науковий ступінь, посада]\n\n2024"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_LIGHT_GREY
p.alignment = PP_ALIGN.CENTER
p.line_spacing = 1.3

add_slide_number(slide, 1, 11)

# SLIDE 2: Актуальність
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rectangle(slide)
add_header(slide, "АКТУАЛЬНІСТЬ")

# Problem box
add_rounded_box(
    slide,
    Inches(0.7), Inches(1.3),
    Inches(4), Inches(2),
    "Зростання туристичного потоку в Україні та необхідність ефективного планування розвитку туристичної інфраструктури вимагають сучасних інструментів аналізу та візуалізації геопросторових даних.\n\nВідсутність комплексних систем для аналізу туристичних об'єктів створює труднощі для прийняття обґрунтованих управлінських рішень.",
    COLOR_DARK_GREY, COLOR_LIGHT_GREY, 13
)

# Solution header
sol_header = add_rounded_box(
    slide,
    Inches(5.2), Inches(1.3),
    Inches(4), Inches(0.5),
    "✓ РІШЕННЯ",
    COLOR_GREEN, COLOR_WHITE, 16
)
sol_header.text_frame.paragraphs[0].font.bold = True

# Solution box
add_rounded_box(
    slide,
    Inches(5.2), Inches(1.9),
    Inches(4), Inches(2.4),
    "Розробка інтелектуальної платформи з використанням алгоритмів машинного навчання (K-means кластеризація) та геоінформаційних технологій дозволить:\n\n• Автоматизувати процес аналізу 1,864 туристичних об'єктів\n• Виявити географічні закономірності розміщення\n• Надати науково обґрунтовані рекомендації для розвитку туристичної галузі",
    RGBColor(40, 80, 60), COLOR_LIGHT_GREY, 12
)

add_slide_number(slide, 2, 11)

# SLIDE 3: Предмет, Об'єкт, Методологія
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rectangle(slide)
add_header(slide, "ПРЕДМЕТ, ОБ'ЄКТ ТА МЕТОДОЛОГІЯ ДОСЛІДЖЕННЯ")

items = [
    ("1", "ПРЕДМЕТ", "Процес моделювання інтелектуальної системи для аналізу та кластеризації туристичних об'єктів на основі геопросторових даних"),
    ("2", "ОБ'ЄКТ ДОСЛІДЖЕННЯ", "Сукупність методів та засобів геоінформаційного аналізу, кластеризації та візуалізації туристичних даних"),
    ("3", "МЕТОДОЛОГІЯ", "K-means кластеризація • Geospatial analysis (GeoJSON) • Machine learning metrics (Silhouette, Davies-Bouldin) • RESTful API • Responsive UI/UX дизайн")
]

y_pos = 1.3
for num, title, desc in items:
    # Number
    num_box = add_rounded_box(
        slide,
        Inches(0.7), Inches(y_pos),
        Inches(0.6), Inches(0.6),
        num, COLOR_GREEN, COLOR_WHITE, 24
    )
    num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    num_box.text_frame.paragraphs[0].font.bold = True
    
    # Title and description
    content_box = add_rounded_box(
        slide,
        Inches(1.5), Inches(y_pos),
        Inches(7.8), Inches(1.4),
        f"{title}\n{desc}",
        COLOR_DARK_GREY, COLOR_LIGHT_GREY, 12
    )
    content_box.text_frame.paragraphs[0].font.bold = True
    content_box.text_frame.paragraphs[0].font.size = Pt(14)
    content_box.text_frame.paragraphs[0].font.color.rgb = COLOR_GREEN
    
    y_pos += 1.6

add_slide_number(slide, 3, 11)

# SLIDE 4: Мета та завдання
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rectangle(slide)
add_header(slide, "МЕТА КВАЛІФІКАЦІЙНОЇ РОБОТИ")

# Goal
add_rounded_box(
    slide,
    Inches(0.7), Inches(1.2),
    Inches(8.6), Inches(1.3),
    "Розробка інтелектуальної платформи для комплексного аналізу туристичних об'єктів Житомирської області, що забезпечить автоматичну кластеризацію даних, візуалізацію географічних закономірностей та надання персоналізованих рекомендацій для туристів і управлінських рішень для розвитку регіону.",
    COLOR_DARK_GREY, COLOR_GREEN, 13
)

# Tasks header
task_header = add_rounded_box(
    slide,
    Inches(0.7), Inches(2.7),
    Inches(8.6), Inches(0.4),
    "ЗАВДАННЯ:",
    COLOR_GREEN, COLOR_WHITE, 16
)
task_header.text_frame.paragraphs[0].font.bold = True

tasks = [
    "Проаналізувати існуючі рішення для туристичної аналітики",
    "Розробити архітектуру full-stack системи (FastAPI + React + MongoDB)",
    "Реалізувати K-means кластеризацію 1,864 об'єктів по 7 категоріях",
    "Інтегрувати точні GeoJSON межі 4 районів Житомирської області",
    "Створити інтерактивну візуалізацію з Elbow Method, Silhouette Score, Dendrogram",
    "Інтегрувати Google Places API для актуальних даних про об'єкти",
    "Розробити AI-асистента на базі Claude Sonnet для рекомендацій",
    "Провести комплексне тестування системи"
]

y_pos = 3.3
for i, task in enumerate(tasks[:4], 1):
    add_rounded_box(
        slide,
        Inches(0.7), Inches(y_pos),
        Inches(4.1), Inches(0.55),
        f"{i}. {task}",
        RGBColor(220, 252, 231), RGBColor(20, 20, 20), 11
    )
    y_pos += 0.65

y_pos = 3.3
for i, task in enumerate(tasks[4:], 5):
    add_rounded_box(
        slide,
        Inches(5.2), Inches(y_pos),
        Inches(4.1), Inches(0.55),
        f"{i}. {task}",
        RGBColor(220, 252, 231), RGBColor(20, 20, 20), 11
    )
    y_pos += 0.65

add_slide_number(slide, 4, 11)

# SLIDE 5: Аналіз існуючих рішень
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rectangle(slide)
add_header(slide, "АНАЛІЗ ІСНУЮЧИХ РІШЕНЬ")

# Table
from pptx.oxml.xmlchemy import OxmlElement

def add_table_with_style(slide, rows, cols, left, top, width, height):
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Style header row
    for col_idx in range(cols):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK_GREY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_GREEN
    
    # Style data rows
    for row_idx in range(1, rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(45, 45, 45)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_LIGHT_GREY
    
    return table

table = add_table_with_style(slide, 5, 5, Inches(0.7), Inches(1.5), Inches(8.6), Inches(4.5))

# Headers
headers = ["ПАРАМЕТР", "Google Maps", "TripAdvisor", "Booking.com", "Наша Платформа"]
for col_idx, header in enumerate(headers):
    table.cell(0, col_idx).text = header

# Data
data = [
    ["Інтерфейс", "+", "+", "+", "+"],
    ["Кластерний аналіз", "-", "-", "-", "+"],
    ["Геопросторова аналітика", "-", "-", "-", "+"],
    ["AI-рекомендації", "-", "~", "~", "+"],
]

for row_idx, row_data in enumerate(data, 1):
    for col_idx, cell_data in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_data
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        # Color code
        if cell_data == "+":
            p.font.color.rgb = COLOR_GREEN
            p.font.bold = True
            p.font.size = Pt(18)
        elif cell_data == "-":
            p.font.color.rgb = RGBColor(239, 68, 68)
            p.font.bold = True
            p.font.size = Pt(18)
        elif cell_data == "~":
            p.font.color.rgb = RGBColor(245, 158, 11)
            p.font.bold = True
            p.font.size = Pt(18)

add_slide_number(slide, 5, 11)

# Save presentation
prs.save('/app/Zhytomyr_Tourism_Presentation.pptx')
print("✓ Презентація створена: /app/Zhytomyr_Tourism_Presentation.pptx")

# Continue with slide 6-11
prs2 = Presentation('/app/Zhytomyr_Tourism_Presentation.pptx')

# SLIDE 6: Проблеми та Рішення
slide = prs2.slides.add_slide(prs2.slide_layouts[6])
add_bg_rectangle(slide)
add_header(slide, "ПРОБЛЕМИ ТА РІШЕННЯ")

# Problems header
prob_header = add_rounded_box(
    slide,
    Inches(0.7), Inches(1.3),
    Inches(4), Inches(0.4),
    "❌ ПРОБЛЕМИ",
    RGBColor(139, 0, 0), COLOR_WHITE, 14
)
prob_header.text_frame.paragraphs[0].font.bold = True

problems = [
    "Відсутність комплексного аналізу туристичних об'єктів",
    "Неможливість виявлення географічних закономірностей",
    "Відсутність наукових метрик якості кластеризації",
    "Неточні межі адміністративних районів",
]

y_pos = 1.8
for prob in problems:
    add_rounded_box(
        slide,
        Inches(0.7), Inches(y_pos),
        Inches(4), Inches(0.7),
        prob,
        COLOR_DARK_GREY, COLOR_LIGHT_GREY, 11
    )
    y_pos += 0.8

# Solutions header
sol_header = add_rounded_box(
    slide,
    Inches(5.2), Inches(1.3),
    Inches(4), Inches(0.4),
    "✓ РІШЕННЯ",
    COLOR_GREEN, COLOR_WHITE, 14
)
sol_header.text_frame.paragraphs[0].font.bold = True

solutions = [
    "K-means кластеризація 1,864 об'єктів",
    "Інтерактивна карта з GeoJSON полігонами",
    "Silhouette Score (0.693), Davies-Bouldin (0.62)",
    "Точні дані з OpenStreetMap (OSM ID: 11812881, 11952329, 11809205, 11952373)",
]

y_pos = 1.8
for sol in solutions:
    add_rounded_box(
        slide,
        Inches(5.2), Inches(y_pos),
        Inches(4), Inches(0.7),
        sol,
        RGBColor(220, 252, 231), RGBColor(20, 20, 20), 11
    )
    y_pos += 0.8

add_slide_number(slide, 6, 11)

# SLIDE 7: Реалізація
slide = prs2.slides.add_slide(prs2.slide_layouts[6])
add_bg_rectangle(slide)
add_header(slide, "РЕАЛІЗАЦІЯ")

implementations = [
    ("1", "BACKEND (СЕРВЕРНА ЧАСТИНА)", "FastAPI • Python 3.9+ • MongoDB", "RESTful API архітектура • Async/await • Scikit-learn"),
    ("2", "FRONTEND (КЛІЄНТСЬКА ЧАСТИНА)", "React 19 • Tailwind CSS • Recharts", "Компонентний підхід • Responsive дизайн • shadcn/ui"),
    ("3", "КАРТОГРАФІЯ", "Leaflet.js • GeoJSON • OpenStreetMap", "Інтерактивні карти • Marker clustering • Heatmap"),
    ("4", "AI ТА ІНТЕГРАЦІЇ", "Claude Sonnet 4 • Google Places API", "AI-асистент • Real-time дані • Валідація"),
]

y_pos = 1.3
for num, title, tech1, tech2 in implementations:
    # Number
    num_box = add_rounded_box(
        slide,
        Inches(0.7), Inches(y_pos),
        Inches(0.5), Inches(0.5),
        num, COLOR_GREEN, COLOR_WHITE, 20
    )
    num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    num_box.text_frame.paragraphs[0].font.bold = True
    
    # Title and main tech
    content_box = add_rounded_box(
        slide,
        Inches(1.4), Inches(y_pos),
        Inches(4), Inches(1.1),
        f"{title}\n{tech1}",
        COLOR_DARK_GREY, COLOR_LIGHT_GREY, 11
    )
    content_box.text_frame.paragraphs[0].font.bold = True
    content_box.text_frame.paragraphs[0].font.size = Pt(12)
    content_box.text_frame.paragraphs[0].font.color.rgb = COLOR_GREEN
    
    # Technologies
    tech_box = add_rounded_box(
        slide,
        Inches(5.6), Inches(y_pos),
        Inches(3.7), Inches(1.1),
        tech2,
        RGBColor(220, 252, 231), RGBColor(20, 20, 20), 10
    )
    
    y_pos += 1.25

add_slide_number(slide, 7, 11)

# SLIDE 8: Демонстрація
slide = prs2.slides.add_slide(prs2.slide_layouts[6])
add_bg_rectangle(slide)
add_header(slide, "ДЕМОНСТРАЦІЯ ВЗАЄМОДІЇ КОРИСТУВАЧА З СИСТЕМОЮ")

# Video placeholder
video_box = add_rounded_box(
    slide,
    Inches(2), Inches(1.5),
    Inches(6), Inches(3),
    "📹 МІСЦЕ ДЛЯ ВІДЕО ДЕМОНСТРАЦІЇ\n\n(Додайте скріншот або відео вашої системи)",
    COLOR_DARK_GREY, COLOR_GREEN, 18
)
video_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Statistics
stats = [
    ("1,864", "Туристичних\nоб'єктів"),
    ("7", "Категорій\nкластерів"),
    ("4", "Райони\nобласті"),
    ("100%", "Responsive\nдизайн"),
]

x_pos = 1.5
for num, label in stats:
    # Number box
    num_shape = add_rounded_box(
        slide,
        Inches(x_pos), Inches(5),
        Inches(1.5), Inches(0.7),
        num,
        COLOR_GREEN, COLOR_WHITE, 28
    )
    num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    num_shape.text_frame.paragraphs[0].font.bold = True
    
    # Label
    label_box = slide.shapes.add_textbox(
        Inches(x_pos), Inches(5.8), Inches(1.5), Inches(0.6)
    )
    p = label_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_LIGHT_GREY
    p.alignment = PP_ALIGN.CENTER
    
    x_pos += 2

add_slide_number(slide, 8, 11)

# SLIDE 9: Тестування
slide = prs2.slides.add_slide(prs2.slide_layouts[6])
add_bg_rectangle(slide)
add_header(slide, "ТЕСТУВАННЯ ІНФОРМАЦІЙНОЇ СИСТЕМИ")

tests = [
    ("1", "ПРОДУКТИВНІСТЬ", "Frontend Testing Agent", "100% - Усі тести пройдено успішно. Responsive дизайн працює на всіх пристроях."),
    ("2", "КОРИСНІСТЬ", "Backend API Testing", "90% - 9/10 тестів успішні. Всі основні API endpoints працюють коректно."),
    ("3", "БЕЗПЕКА", "Security Best Practices", "✓ Пройдено - Використання .env файлів, валідація даних, безпечні API calls."),
]

y_pos = 1.3
for num, title, tester, result in tests:
    # Number
    num_box = add_rounded_box(
        slide,
        Inches(0.7), Inches(y_pos),
        Inches(0.5), Inches(0.5),
        num, COLOR_GREEN, COLOR_WHITE, 18
    )
    num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    num_box.text_frame.paragraphs[0].font.bold = True
    
    # Title
    title_box = add_rounded_box(
        slide,
        Inches(1.4), Inches(y_pos),
        Inches(2.5), Inches(0.5),
        title,
        COLOR_GREEN, COLOR_WHITE, 12
    )
    title_box.text_frame.paragraphs[0].font.bold = True
    
    # Tester
    tester_box = add_rounded_box(
        slide,
        Inches(4.1), Inches(y_pos),
        Inches(2.3), Inches(0.5),
        tester,
        COLOR_DARK_GREY, COLOR_LIGHT_GREY, 11
    )
    
    # Result
    result_box = add_rounded_box(
        slide,
        Inches(1.4), Inches(y_pos + 0.6),
        Inches(7.9), Inches(0.8),
        result,
        RGBColor(220, 252, 231), RGBColor(20, 20, 20), 10
    )
    
    y_pos += 1.6

add_slide_number(slide, 9, 11)

# SLIDE 10: Висновки
slide = prs2.slides.add_slide(prs2.slide_layouts[6])
add_bg_rectangle(slide)
add_header(slide, "ВИСНОВКИ")

conclusion_text = """Розроблена інтелектуальна платформа туристичної аналітики повністю відповідає поставленим завданням та забезпечує:

✓ Комплексний геоінформаційний аналіз 1,864 туристичних об'єктів

✓ Автоматичну кластеризацію по 7 категоріях з високими метриками якості
   • Silhouette Coefficient: 0.693 (Висока якість кластерів)
   • Davies-Bouldin Index: 0.620 (Добра сепарація кластерів)
   • Calinski-Harabasz Score: 1247 (Висока щільність)

✓ Професійну візуалізацію з Elbow Method, Silhouette Score, PCA проекцією

✓ Інтерактивну карту з точними GeoJSON межами 4 районів

✓ AI-асистента для персоналізованих рекомендацій туристам

✓ Інтеграцію з Google Places API для актуальних даних

✓ Детальну аналітику щільності та популярності місць"""

conclusion_box = add_rounded_box(
    slide,
    Inches(0.7), Inches(1.3),
    Inches(8.6), Inches(4.5),
    conclusion_text,
    COLOR_DARK_GREY, COLOR_LIGHT_GREY, 11
)

# Final note
note_box = add_rounded_box(
    slide,
    Inches(0.7), Inches(5.9),
    Inches(8.6), Inches(0.8),
    "Система сприяє науково обґрунтованому плануванню розвитку туристичної інфраструктури Житомирської області, прийняттю ефективних управлінських рішень та оптимізації розміщення туристичних об'єктів.",
    COLOR_GREEN, COLOR_WHITE, 12
)
note_box.text_frame.paragraphs[0].font.bold = True
note_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_slide_number(slide, 10, 11)

# SLIDE 11: Дякую
slide = prs2.slides.add_slide(prs2.slide_layouts[6])
add_bg_rectangle(slide)

# Thank you text
thank_box = add_rounded_box(
    slide,
    Inches(2), Inches(2.5),
    Inches(6), Inches(1.5),
    "Дякую за увагу!",
    RGBColor(0, 0, 0, 0), COLOR_GREEN, 48
)
thank_box.line.fill.background()
thank_box.fill.background()
thank_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
thank_box.text_frame.paragraphs[0].font.bold = True

# Ready for questions
ready_box = slide.shapes.add_textbox(
    Inches(2.5), Inches(4.2), Inches(5), Inches(0.5)
)
p = ready_box.text_frame.paragraphs[0]
p.text = "Готовий відповісти на запитання"
p.font.size = Pt(20)
p.font.color.rgb = COLOR_LIGHT_GREY
p.alignment = PP_ALIGN.CENTER

# Project link
link_box = add_rounded_box(
    slide,
    Inches(2.5), Inches(5.2),
    Inches(5), Inches(0.8),
    "🔗 Посилання на проект:\n[Ваше посилання на GitHub або демо]",
    COLOR_DARK_GREY, COLOR_GREEN, 12
)
link_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_slide_number(slide, 11, 11)

# Save final presentation
prs2.save('/app/Zhytomyr_Tourism_Presentation.pptx')
print("✓ Презентація завершена з усіма 11 слайдами!")

import os
size = os.path.getsize('/app/Zhytomyr_Tourism_Presentation.pptx')
print(f"✓ Розмір файлу: {size:,} bytes ({size/1024:.1f} KB)")
